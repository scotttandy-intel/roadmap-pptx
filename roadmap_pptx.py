"""
roadmap_pptx.py

Objects related to the creation of roadmap slides in PowerPoint

Author Scott Tandy
Copyright Intel Corporation 2020

"""

from collections import OrderedDict
import copy

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.fill import FillFormat
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.shapes.connector import Connector
from pptx.dml.line import LineFormat
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.shapes.connector import Connector    ### Used for a patch

import intel_ww as iw
import roadmap_table as rt
import annotations as an


### Patch for pptx
def get_or_add_ln(self):
    return self._element.spPr.get_or_add_ln()

Connector.get_or_add_ln = get_or_add_ln



def toCm(n):
    return n / float(Cm(n))

NAME = 'Name'
OD = 'OrderedDict'

"""
PRODUCT SHAPE
"""
PRODUCT_SHAPE_HEIGHT_CM = 0.4
PRODUCT_TEXT_IN_SHAPE_PTS = Pt(6.0)
PRODUCT_SHAPE_LINE_WIDTH = Pt(1.0)

"""
MILESTONES
"""
PRODUCT_MILESTONE_SHAPE_HEIGHT_CM = 0.32
PRODUCT_MILESTONE_TEXT_IN_SHAPE_PTS = Pt(6.0)
ADD_TO_CENTER_MILESTONE_IN_PRODUCT_SHAPE = 0.05 +\
    (PRODUCT_MILESTONE_SHAPE_HEIGHT_CM - PRODUCT_MILESTONE_SHAPE_HEIGHT_CM) / 2.0

"""
GOLDEN ROADMAP LAYOUT
"""
ROADMAP_TOP = 4.5
ROADMAP_ZERO_LEFT = 2.0
#ROADMAP_HEIGHT = 12.7
ROADMAP_HEIGHT = 13.5
ROADMAP_WIDTH = 31.0
NAME_LEFT_EDGE = 0.5
YEAR_ZERO = 19
QUARTER_SPAN = 20

WW_PER_Q = 13
POINTS_PER_CM = 0.035
SM_CM = 0.2

DEFAULT_FONT_NAME = 'Intel Clear'
DEFAULT_FONT_SIZE_PT = 10

TWO_MM = 0.2


def rotate_table_cell( table_cell, rotation=270 ):
    """
    Some python_pptx magic code from Stackoverflow....
    https://stackoverflow.com/questions/54828103/how-to-rotate-text-in-the-table-cell

    Looks like vert90 is also a likely valid value
    """
    VERT = 'vert'

    tcpr = table_cell._tc.get_or_add_tcPr()

    tcpr.set(VERT,VERT+str(rotation)) # setting 'vert' parameter to 'vert270'

def add_golden_roadmap_slide( presentation ):
    print('add_golden_roadmap_slide()')
    GOLDEN_LAYOUT_NAME = 'Right Quarter white new'

    layout_dict = {}
    for layout in presentation.slide_layouts:
        layout_dict[layout.name]=layout
    
    golden_slide = presentation.slides.add_slide(layout_dict[GOLDEN_LAYOUT_NAME])

    print(golden_slide.slide_layout.name)

    return golden_slide

def text_length_to_cm( text, pts ):
    len_cm = len(text) * pts * POINTS_PER_CM

    if text == 'A0 TI':
        len_cm = len_cm * 0.8
        
    return len_cm

class Milestone:
    def __init__(self, text, ww_date, adjacent_text ='', \
        fill_rgb=RGBColor(0,0,0), text_rgb=RGBColor(255,255,255), \
            text_size=Pt(6), category='Unspecified',\
                marker_left_cm=None, marker_width_cm=None, \
                    text_box_left_cm=None, text_box_width_cm=None):
        self.text = text
        self.ww_date = ww_date    ## of the type rt.WW
        self.adjacent_text = adjacent_text
        self.fill_rgb = fill_rgb
        self.text_rgb = text_rgb
        self.category = category

        ### Values filled in during the rendering process
        self.marker_left_cm = marker_left_cm
        self.marker_width_cm = marker_width_cm
        self.text_box_left_cm = text_box_left_cm
        self.text_box_width_cm = text_box_width_cm
        self.is_first = None
        self.is_last = None

        return

    def __str__(self):
        return str(self.text)+'-'+str(self.ww_date)
    
    def __eq__( self, other):
        return self.ww_date == other.ww_date

    def __lt__(self, other):
        return (self.ww_date < other.ww_date)
    
    def ww(self):
        return self.ww_date.ww
    
    def year(self):
        return self.ww_date.year
    
    def right_of_marker_and_text(self):
        """
        """
        if self.text_box_left_cm is None or self.text_box_width_cm is None:
            right_edge = self.marker_left_cm+self.marker_width_cm
        else:
            if self.marker_left_cm < self.text_box_left_cm:
                right_edge = self.text_box_left_cm +self.text_box_width_cm
            else:
                right_edge = self.marker_left_cm +self.marker_width_cm
        
        return right_edge
    
    def between_two_milestones(self, other_ms):
        """
        This is a very specific function that figures out how much
        space there is between two milestones that were rendered.
        Since they were rendered, the left/width cm parameters are updated.
        {marker}{text}   (space inbetween)  {text}{marker}
        The function doesn't assume anything about the positioning of
        the relative text and markers in a milestone

        returns a tulple (left_edge_cm, right_edge_cm)
        marker_left_cm,marker_width_cm, text_box_left_cm, text_box_width_cm
        """
        ### Which way are the milstones ordered?
        if self.ww_date < other_ms.ww_date:
            left_ms = self
            right_ms = other_ms
        else:
            left_ms = other_ms
            right_ms = self
        
        print('---should be in date order', left_ms.ww_date, right_ms.ww_date)

        #### Within the left milestone, what order are the marker and text?
        ### The assumption is that the milestone was always rendered even if
        ### the text was not...maybe bad assumption....
        if left_ms.text_box_left_cm is None or left_ms.text_box_width_cm is None:
            left_edge_cm = left_ms.marker_left_cm+left_ms.marker_width_cm
        else:
            if left_ms.marker_left_cm < left_ms.text_box_left_cm:
                left_edge_cm = left_ms.text_box_left_cm +left_ms.text_box_width_cm
            else:
                left_edge_cm = left_ms.marker_left_cm +left_ms.marker_width_cm
        
        #### Within the right milestone, what order are the marker and text?
        if right_ms.text_box_left_cm is None or right_ms.text_box_width_cm is None:
            right_edge_cm = right_ms.marker_left_cm
        else:
            if right_ms.marker_left_cm < right_ms.text_box_left_cm:
                right_edge_cm = right_ms.marker_left_cm
            else:
                right_edge_cm = right_ms.text_box_left_cm
        
        return (left_edge_cm, right_edge_cm)

class Component:
    def __init__(self, component_name, category=None):
        self.component_name = component_name
        self.category = category
        return

class ProgramInformation:
    def __init__(self, shorthand_name, milestones_list, from_roadmap_top_cm, annotations_list=None,
            component_list=None):
        self.shorthand_name = shorthand_name
        self.milestones = milestones_list   ## Of type list(Milestones)
        self.from_roadmap_top_cm = from_roadmap_top_cm
        self.component_list = component_list
        return
    
    def __str__(self):
        s = f"{self.shorthand_name} - {self.from_roadmap_top_cm:0.4f}"
        return(s)
    
    def just_ww_milestones(self):
        just_wws = []
        for ms in self.milestones:
            if isinstance( ms.ww_date, iw.WW):
                just_wws.append(ms)
            else:
                print('Skipping MS', ms.ww_date)
        
        return just_wws

    def first_milestone(self):
        return self.just_ww_milestones()[0]
    
    def last_milestone(self):
        return self.just_ww_milestones()[-1]
    
    def milestone_by_name(self, milestone_name):
        for ms in self.milestones:
            if ms.text.upper() == milestone_name.upper():
                return ms
        return None
    
    def update_milestones(self, new_milestone_list):
        """
        Updates the milestone list - still requires the 
        caller to know a bit about the strucure of this class
        """
        self.milestones = new_milestone_list

class ConceptProgramInformation:
    def __init__(self, si_program, from_roadmap_top_cm ):
        self.si_program = si_program
        self.shorthand_name = self.si_program.name
        self.milestones = self._create_milestones_list()
        self.from_roadmap_top_cm = from_roadmap_top_cm
        return
    
    def __str__(self):
        s = f"{self.shorthand_name} - {self.from_roadmap_top_cm:0.4f}"
        return(s)
    
    def _create_milestones_list(self):
        ms_dict  = self.si_program.milestones()

        ms_list = []

        for key, value in ms_dict.items():
            ms_list.append(Milestone(text=key, adjacent_text=str(value), ww_date=value))
        
        if len(ms_list) == 0:
            raise ValueError("No milestones found for: " + self.si_program.name)
        
        if len(ms_list) < 2:
            print("Adding milestone for ", self.si_program.name)
            earlier_ww = ms_list[0].ww_date.add_wws(wws=26)
            ms_list.append(Milestone(text='????', ww_date=earlier_ww))
        
        ms_list = sorted(ms_list)
        
        return ms_list

    def first_milestone(self):
        return self.milestones[0]
    
    def last_milestone(self):
        return self.milestones[-1]

    def milestone_before(self, col_name):
        ms_name = self.si_program.milestone_before(col_name=col_name)
        return self.milestone_by_text(milestone_name=ms_name)
    
    def milestone_after(self, col_name):
        ms_name = self.si_program.milestone_after(col_name=col_name)
        return self.milestone_by_text(milestone_name=ms_name)
    
    def update_milestones(self, new_milestone_list):
        """
        Updates the milestone list - still requires the 
        caller to know a bit about the strucure of this class
        """
        self.milestones = new_milestone_list
    
    def milestone_by_text(self, milestone_name):
        for ms in self.milestones:
            if ms.text == milestone_name:
                return ms
        return None

    def annotations(self):
        return self.si_program.annotations()

def rt_to_ProgramInformation(roadmap_table, roadmap_configuration, shorthand_name, from_roadmap_top_cm, wws_back=26, missing_ms='????'):
    
    sp = rt.SiProduct(shorthand_name=shorthand_name, roadmap_table=roadmap_table)

    ms_dict  = sp.milestone_dict(tag_list=roadmap_configuration.milestones.values())

    ms_list = []

    for key, value in ms_dict.items():
        ms_list.append(Milestone(text=key, adjacent_text=str(value), ww_date=value))
    
    if len(ms_list) < 2:
        print("Adding milestone for ", sp.name)
        earlier_ww = ms_list[0].ww_date.subtract_wws_return_new(wws=(26))
        ms_list.append(Milestone(text=missing_ms, ww_date=earlier_ww))
    
    ms_list = sorted(ms_list)

    pi = ProgramInformation(shorthand_name=sp.shorthand_name, from_roadmap_top_cm=from_roadmap_top_cm, milestones_list=ms_list)

    return pi

def concept_rt_to_ProgramInformation(major_value, minor_value, name, my_row, roadmap_configuration, \
    from_roadmap_top_cm, wws_back=26, missing_ms='????'):
    
    sp = rt.ConceptSiProduct(major_value=major_value, minor_value=minor_value, name=name, my_row=my_row)

    pi = ConceptProgramInformation(si_program=sp, from_roadmap_top_cm=from_roadmap_top_cm)

    return pi

class ProgramComponents:
    def __init__(self, family_name, program_list=None):
        self.family_name = family_name
        self.program_list = program_list
        return

class Segment:
    def __init__(self, segment_name, program_list=None):
        self.segment_name = segment_name
        self.program_list = program_list
        return

class Business:
    def __init__(self, business_name, segment_list=None):
        self.business_name = business_name
        self.program_list = segment_list
        return

class CanvasGrid:
    """
    Class to collect up helper functions to figure out how to draw the 
    roadmap year header and quarter header.
    """
    def __init__(self, start_ww, end_ww ):
        self.start_ww = start_ww
        self.start_qtr = start_ww.qtr_of_year()
        self.start_year = start_ww.year
        self.left_edge_ww = iw.WW(ww=(self.start_qtr-1)*WW_PER_Q+1, year=self.start_year)

        self.end_ww = end_ww
        self.end_qtr = end_ww.qtr_of_year()
        self.end_year = end_ww.year
        #self.right_edge_ww = iw.WW(ww=(self.end_qtr)*WW_PER_Q,year=self.end_year)
        self.right_edge_ww = iw.WW(ww=(self.end_qtr)*WW_PER_Q,year=self.end_year)

        self.ww_span = self.left_edge_ww.ww_delta_from(self.right_edge_ww)

        self.ww_to_pct = build_canvas_dict(left_ww=self.left_edge_ww,\
             right_ww=self.right_edge_ww, quarter_count=self.total_quarters())
        return
    
    def __str__(self):
        return 'From:'+str(self.start_ww) + ' to ' +'To:'+str(self.end_ww)
    
    def year_span_count(self):
        return self.end_year - self.start_year + 1
    
    def quarters_in_first_year(self):
        """
        How many quaters in the first year do we need to show?
        """
        if self.year_span_count() == 1:
            return self.end_qtr - self.start_qtr + 1

        return iw.QTR_PER_YEAR -  self.start_qtr + 1
    
    def quarters_in_last_year(self):
        if self.year_span_count() == 1:
            raise ValueError('Quaters in last year not defined if you only span 1 year!')
        
        return self.end_qtr

    def total_quarters(self):
        if self.year_span_count() == 1:
            return self.quarters_in_first_year()
        fyq = self.quarters_in_first_year()
        myq = (self.year_span_count() - 2) * iw.QTR_PER_YEAR
        lyq = self.quarters_in_last_year()

        return  fyq + myq + lyq
    
    def ww_to_pct_width(self, ww):
        return self.ww_to_pct[dict_str(ww=ww.ww, year=ww.year)]

def dict_str(ww, year):
    return f"{ww:02}'{year:02}"

def build_canvas_dict(left_ww, right_ww, quarter_count):
    """
    Maps ww's to the X axis.
    """
    ww_to_pct = OrderedDict()

    one_ww_pct = 1.0 / (quarter_count * iw.WW_PER_QTR)
    
    start_year = left_ww.year
    start_ww = left_ww.ww

    end_year = right_ww.year
    end_ww = right_ww.ww
    pct_sum = 0.0

    for year in range(start_year, end_year+1):
        ### This is the tricky bit
        ### and is a compromise since it spreads out all
        ### 53 ww's throughout the year instead of cramming
        ### an extra WW in Q4 - shouldn't be too noticeable
        ### and avoides cumulative error across years
        ww_this_year = iw.wws_in_year(year+2000)
        my_ww_pct = one_ww_pct * (iw.WW_PER_YEAR_USUALLY / ww_this_year)
        for ww in range(1, ww_this_year+1):
            ### Skip the weeks at the start
            if ww < start_ww:
                continue
            
            ww_to_pct[dict_str(ww=ww, year=year)] = pct_sum
    
            ### Skip the ww's at the end
            if year == end_year and ww == end_ww:
                break
            
            pct_sum =pct_sum + my_ww_pct
        start_ww = 1
    
    return ww_to_pct

class RoadmapCanvas:
    DEFAULT_YEAR_HEIGHT_CM = 0.75
    DEFAULT_QUARTER_HEIGHT_CM = 0.55

    def __init__(self, grid, canvas_left_edge_cm = 4.0, canvas_top_cm=3, \
        canvas_height_cm=13.6, canvas_width_cm=28.5, years_height_cm=0.75,\
            qtrs_height_cm=0.6 ):
        self.grid = grid
        self.canvas_left_edge_cm = canvas_left_edge_cm
        self.canvas_top_cm = canvas_top_cm
        self.canvas_height_cm = canvas_height_cm
        self.canvas_width_cm = canvas_width_cm
        self.years_height_cm = years_height_cm
        self.qtrs_height_cm = qtrs_height_cm
        self.draw_area_height = self.canvas_height_cm - \
            (self.years_height_cm + self.qtrs_height_cm)
        self.top_of_draw_area = self.canvas_top_cm + \
            (self.years_height_cm + self.qtrs_height_cm)
        
        self.cm_per_ww = self.canvas_width_cm / float(self.grid.total_quarters() * WW_PER_Q)
        return
    
    def ww_to_slide_x_cm(self, ww):
        pct_of_x_canvas = self.grid.ww_to_pct_width(ww=ww)
        from_left_edge_cm = pct_of_x_canvas * self.canvas_width_cm
        return self.canvas_left_edge_cm + from_left_edge_cm
    
    def test_circle(self, shapes, abs_x_cm, abs_y_cm=None, text='test', width_cm=0.10,\
        font_size=Pt(4) ):
        if abs_y_cm is None:
            abs_y_cm=self.top_of_draw_area
        
        add_rounded_rectangle_marker(shapes=shapes, left_cm=abs_x_cm,\
             top_cm=abs_y_cm, width_cm=width_cm, height_cm=PRODUCT_SHAPE_HEIGHT_CM, \
                 text = text, font_name='Intel Clear', font_size = font_size, )

    def render_yrs_qts_table(self, shapes, align_zero=False ):
        """
        Renders YrQtWWGrid object into slide
        """
        frame = shapes.add_table(rows=3,\
            cols=self.grid.total_quarters(),\
            left=Cm(self.canvas_left_edge_cm),\
                top=Cm(self.canvas_top_cm), \
            width=Cm(self.canvas_width_cm),\
            height=Cm(self.canvas_height_cm) )

        table = frame.table

        # Defaults from the  Golden Roadmap
        table.first_col = False
        table.first_row = True
        table.horz_banding = True
        table.last_col = False
        table.last_row = False
        table.vert_banding = False

        # Set the row heights for the headers and draw area
        table.rows[0].height = Cm(self.years_height_cm)
        table.rows[1].height = Cm(self.qtrs_height_cm)
        table.rows[2].height = Cm(self.draw_area_height)
        
        qts_each_year = [self.grid.quarters_in_first_year()]

        if not align_zero:
            year_text = [f"20{self.grid.start_year:02d}"]
        else:
            year_text = [f"{self.grid.start_year:02d}"] 

        q_start = [self.grid.start_qtr]

        if (self.grid.year_span_count() > 1):
            ## One year - may only be a subset of the quarters
            for y in range(self.grid.year_span_count() - 2):
                qts_each_year.append(iw.QTR_PER_YEAR)
                year_text.append(f"20{self.grid.start_year+y+1:02d}")
                q_start.append(1)
            
            qts_each_year.append(self.grid.quarters_in_last_year())
            if not align_zero:
                year_text.append(f"20{self.grid.end_year:02d}")
            else:
                year_text.append(f"{self.grid.end_year:02d}")
            q_start.append(1)

        current_col = 0
        row_zero_cells = table.rows[0].cells
        
        for merge_qts, year_text in zip(qts_each_year,year_text):
            row_zero_cells[current_col].text = year_text
            row_zero_cells[current_col].merge(row_zero_cells[current_col+merge_qts-1])
            font_size = 14
            
            if merge_qts < 2: ## Shrink font size if we ar squeesing into just one Q
                font_size = 10
            
            self.format_table_cell_centered(cell=row_zero_cells[current_col],\
                font_size=font_size)
            
            current_col = current_col + merge_qts
        current_col = 0
        qtr_cells = table.rows[1].cells
        for qs,merge_qts in zip(q_start, qts_each_year):
            for i in range(merge_qts):
                qtr_cells[current_col+i].text = f"Q{i+qs}"
                self.format_table_cell_centered(cell=qtr_cells[current_col+i],font_size=11)
            current_col = current_col + merge_qts
        
        return

    def format_table_cell_centered(self, cell, font_size = DEFAULT_FONT_SIZE_PT, font_name=DEFAULT_FONT_NAME ):
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Cm(0.0)
        cell.margin_right = Cm(0.0)
        cell.margin_top = Cm(0.0)
        cell.margin_bottom = Cm(0.0)

        tf = cell.text_frame
        tf.auto_size = None
        tf.margin_bottom = Cm(0.1)
        tf.margin_left = Cm(0.1)
        tf.margin_right = Cm(0.1)
        tf.margin_top = Cm(0.1)
        tf.word_wrap = None
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        para.font.name = font_name
        para.font.size = Pt(font_size)
        para.font.bold = True

        return


class SwimlanesGrid:
    ROWS = 'ROWS'
    ROW_COUNT = 'RowCount'
    ROW_INDEX = 'RowNumber'
    NAME = 'Name'

    def __init__(self, hierarchy, major_column_name, minor_column_name ):
        """
        herarchy is in the form:
        h = OrderedDict()
        h['Datacenter']=OrderedDict()
        h['Datacenter']['Compute']='Compute'
        h['Datacenter']['ComputePC']=''
        h['Datacenter']['MultiPurpose']='Multi Purpose'
        h['Discrete Client']=OrderedDict()
        h['Discrete Client']['Pro']='Pro'
        h['Discrete Client']['Enthusiast']='Enthusiast'
        h['Discrete Client']['Performance']='Performance'
        h['Discrete Client']['Mainstream']='Mainstream'
        h['Discrete Client']['Entry']='Entry'

        NOTE: This function uses a convention that if you provde ''
        as a name for a segment - this is used only to create
        an additional row of space in the layout - and is not
        included in the reporting of actual segments in the table
        """
        self.h = OrderedDict()
        self.row_count = 0
        for major, major_d in hierarchy.items():
            self.h[major] = OrderedDict()
            self.h[major][SwimlanesGrid.ROW_COUNT]=len(major_d)
            self.h[major][SwimlanesGrid.ROWS]=OrderedDict()
            for minor, minor_namme in major_d.items():
                self.h[major][SwimlanesGrid.ROWS][minor]=OrderedDict()
                self.h[major][SwimlanesGrid.ROWS][minor][SwimlanesGrid.NAME]=minor_namme
                self.h[major][SwimlanesGrid.ROWS][minor][SwimlanesGrid.ROW_INDEX]=self.row_count
                self.row_count = self.row_count + 1
        return
    
    def major_row_values(self):
        return list(self.h.keys())
    
    def minor_row_values(self, major_row_value):
        minor_row_list = []
        for segment in self.h[major_row_value][SwimlanesGrid.ROWS].keys():
            if self.h[major_row_value][SwimlanesGrid.ROWS][segment][SwimlanesGrid.NAME] == '':
                continue
            minor_row_list.append(segment)
        return minor_row_list
    
    def rows_for_major_value(self, major_value):
        """
        Used in figuring out how many cells to merge for the business
        """
        return self.h[major_value][SwimlanesGrid.ROW_COUNT]
    
    def row_index_for_business_segment(self, business_name, segment_name):
        return self.h[business_name][SwimlanesGrid.ROWS][segment_name][SwimlanesGrid.ROW_INDEX]
    
    def row_index_for_major_column_minor_column(self, major_value, minor_value):
        return self.h[major_value][SwimlanesGrid.ROWS][minor_value][SwimlanesGrid.ROW_INDEX]
    
    def segment_name_for_business_segment(self, business_name, segment_name):
        return self.h[business_name][SwimlanesGrid.ROWS][segment_name][SwimlanesGrid.NAME]
    
    def third_column_for_major_column_minor_column(self, major_value, minor_value):
        return self.h[major_value][SwimlanesGrid.ROWS][minor_value][SwimlanesGrid.NAME]

class SwimlaneTable:
    """
    Function to generate the Swimlanes Table
    """
    def __init__(self, grid, table_left_edge_cm = 2.0, table_top_cm=4.0, \
        table_height_cm=12.0, table_width_cm=5.5 ):
        
        self.grid = grid
        self.table_left_edge_cm = table_left_edge_cm
        self.table_top_cm = table_top_cm
        self.table_width_cm = table_width_cm
        self.table_height_cm = table_height_cm
        self.row_count = self.grid.row_count
        self.cm_per_row = self.table_height_cm / float(self.row_count)

        return
    
    def y_value_in_cm_for_segment( self, business, segment):
        row_index = self.grid.row_index_for_business_segment(business_name=business,\
            segment_name=segment)
        
        return self.table_top_cm + (row_index * self.cm_per_row)

    def y_value_in_cm_for_major_column_minor_column( self, major_value, minor_value ):
        row_index = \
            self.grid.row_index_for_major_column_minor_column(major_value=major_value, \
                minor_value=minor_value)
        
        return self.table_top_cm + (row_index * self.cm_per_row)
    
    def render_swimlanes_table(self, shapes ):
            """
            Renders SwimlanesGrid into table
            """
            frame = shapes.add_table(rows=self.grid.row_count,\
                cols=2,\
                left=Cm(self.table_left_edge_cm),\
                top=Cm(self.table_top_cm), \
                width=Cm(self.table_width_cm),\
                height=Cm(self.table_height_cm) )

            table = frame.table

            ## Set column widths
            table.columns[0].width = Cm(0.9)

            # Defaults from the  Golden Roadmap
            table.first_col = False
            table.first_row = False
            table.horz_banding = False
            table.last_col = False
            table.last_row = False
            table.vert_banding = False
            
            row_count = 0
            
            for major in self.grid.major_row_values():
                major_row_cnt = self.grid.rows_for_major_value(major)
                table.cell(row_count, 0).text = major
                
                self.format_business_table_cell(table.cell(row_count, 0),\
                    font_size=12)
                
                table.cell(row_count, 0).merge(table.cell(row_count+major_row_cnt-1,0))

                row_count = row_count + self.grid.rows_for_major_value(major) 

                for minor in self.grid.minor_row_values(major_row_value=major):
                    row_index = \
                        self.grid.row_index_for_major_column_minor_column(major_value=major, minor_value=minor)
                    
                    segment_name = \
                        self.grid.third_column_for_major_column_minor_column(major_value=major, minor_value=minor)
                    
                    table.cell(row_index, 1).text = segment_name
                    
                    self.format_segment_table_cell(table.cell(row_index, 1),\
                        font_size=10)
            return
    
    def format_business_table_cell(self, cell, font_size = DEFAULT_FONT_SIZE_PT,\
         font_name=DEFAULT_FONT_NAME ):

            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Cm(0.0)
            cell.margin_right = Cm(0.0)
            cell.margin_top = Cm(0.0)
            cell.margin_bottom = Cm(0.0)    

            tf = cell.text_frame
            tf.auto_size = None
            tf.margin_bottom = Cm(0.05)
            tf.margin_left = Cm(0.05)
            tf.margin_right = Cm(0.05)
            tf.margin_top = Cm(0.05)
            tf.word_wrap = None
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

            
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            para.font.name = font_name
            para.font.size = Pt(font_size)
            para.font.bold = True
            
            rotate_table_cell(cell)

            return
    
    def format_segment_table_cell(self, cell, font_size = DEFAULT_FONT_SIZE_PT,\
         font_name=DEFAULT_FONT_NAME ):
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Cm(0.0)
            cell.margin_right = Cm(0.0)
            cell.margin_top = Cm(0.0)
            cell.margin_bottom = Cm(0.0)       

            tf = cell.text_frame
            tf.auto_size = None
            tf.margin_bottom = Cm(0.05)
            tf.margin_left = Cm(0.05)
            tf.margin_right = Cm(0.05)
            tf.margin_top = Cm(0.05)
            tf.word_wrap = None
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            para.font.name = font_name
            para.font.size = Pt(font_size)
            para.font.bold = True
            return

class RoadmapSlide:
    """
    Here is the higherarchy that the RoadmapSlide utilises

    TopLevel: [IP, SOC Product, Board Product, Reference System, Software Product]
    Business:  
        SOC Product: [DataCenter, DataCenter GPU, Client, Client GPU, Integrated]
    
    Segment: 
        SOC Product:Client GPU:[Entry, Mainstream, Performance, Enthusiast, Pro]
        SOC Product:Data Center GPU: [Multi-Use, Compute]

    """
    def __init__(self, slide, roadmap_canvas, roadmap_left_edge_cm = 2.0, roadmap_top_cm=4.5, roadmap_bottom_cm=18.2,\
         roadmap_width_cm=60.0, name_left_edge_cm=1.25, year_zero = YEAR_ZERO, quarter_span=20):
        self.slide = slide
        self.roadmap_canvas = roadmap_canvas
        self.shapes = slide.shapes
        self.roadmap_left_edge_cm = roadmap_left_edge_cm
        self.roadmap_top_cm = roadmap_top_cm
        self.roadmap_bottom_cm =roadmap_bottom_cm
        self.roadmap_width_cm = roadmap_width_cm
        self.name_left_edge_cm = name_left_edge_cm
        self.year_zero = year_zero
        self.quarter_span = quarter_span

        self.cm_per_ww = self.roadmap_width_cm / (float(self.quarter_span) * WW_PER_Q)
        #print('========>',self.cm_per_ww, '=', self.roadmap_width_cm,'/', float(self.quarter_span),'*', WW_PER_Q)
    
    def __str__(self):
        s = ''
        for attr in dir(self):
            if '__' in attr:
                continue
            s = s + ' ' + attr +':' + str(getattr(self, attr))
        return(s)
    
    def ww_to_offset_cm(self, ww, align_zero_ww=None):
        if align_zero_ww is not None:
            ww_delta = ww.ww_delta_from(align_zero_ww)
            if ww_delta < 0:
                ww = ww.add_wws(ww_delta)
            elif ww_delta > 0:
                ww = ww.subtract_wws(-1* ww_delta)

        return self.roadmap_canvas.ww_to_slide_x_cm(ww)
        

    def render_program(self, program_information, from_roadmap_top_cm=None, \
        align_zero=False):
        """
        """
        if from_roadmap_top_cm is None:
            from_roadmap_top_cm = program_information.from_roadmap_top_cm
        
        vertical_offset_cm = from_roadmap_top_cm

        first_milestone = program_information.first_milestone()
        
        if align_zero:
            align_zero_ww = first_milestone.ww_date
        else:
            align_zero_ww = None

        left_offset_cm = self.ww_to_offset_cm(ww=first_milestone.ww_date, \
            align_zero_ww=align_zero_ww)
        
        last_milestone = program_information.last_milestone()

        width_cm = first_milestone.ww_date.ww_delta_from(last_milestone.ww_date) * self.roadmap_canvas.cm_per_ww

        add_rounded_rectangle_marker(shapes=self.slide.shapes, \
            left_cm=left_offset_cm,\
            width_cm = width_cm, top_cm=vertical_offset_cm, text = None,\
                font_size=PRODUCT_TEXT_IN_SHAPE_PTS, font_bold=True,\
                     line_rgb=RGBColor(0,0,0), line_width=PRODUCT_SHAPE_LINE_WIDTH)
        
        prog_left_offset_cm = left_offset_cm - \
            text_length_to_cm(program_information.shorthand_name, PRODUCT_MILESTONE_TEXT_IN_SHAPE_PTS/Pt(1.0))-\
                TWO_MM

        
        text_width_cm = text_length_to_cm(program_information.shorthand_name, PRODUCT_MILESTONE_TEXT_IN_SHAPE_PTS/Pt(1.0))
        
        add_text_box(shapes=self.slide.shapes, left_cm=prog_left_offset_cm, \
            top_cm=vertical_offset_cm+PRODUCT_MILESTONE_SHAPE_HEIGHT_CM/3.0, 
            width_cm = text_width_cm, \
                height_cm=PRODUCT_MILESTONE_SHAPE_HEIGHT_CM,\
                text = program_information.shorthand_name, text_align=PP_ALIGN.RIGHT, text_auto_size=MSO_AUTO_SIZE.NONE,\
                fill_rgb=None, line_rgb = None, line_width= Pt(1.0), line_dash = None,\
                font_name='Intel Clear', font_size = Pt(8), font_rgb_color=RGBColor(0,0,0), font_bold=True, font_italic=False,\
                font_underline=False, font_language_id=MSO_LANGUAGE_ID.NONE)

        vertical_offset_cm = vertical_offset_cm + ADD_TO_CENTER_MILESTONE_IN_PRODUCT_SHAPE
        
        updated_milestone_list = []   ### Make new copy to avoid changing items durin
                                    ### the enumeration below
        for i, ms in enumerate(program_information.milestones):
            is_first, is_last = (False,False)
            if i == 0:
                is_first = True
            elif i == len(program_information.milestones)-1:
                is_last = True
            
            updated_milestone = self.render_milestone(milestone=ms, vertical_offset_cm=vertical_offset_cm,\
                 is_first=is_first, is_last=is_last)

            updated_milestone_list.append(updated_milestone)
        
        ### Yes we need to know too much about the program_information object to
        program_information.update_milestones(updated_milestone_list)

        return program_information
    
    def render_milestone(self, milestone, vertical_offset_cm, left_offset_cm=0.0, \
        is_first=False, is_last=False):
        
        shape_text = milestone.text
        width_cm = text_length_to_cm(shape_text, PRODUCT_MILESTONE_TEXT_IN_SHAPE_PTS/Pt(1.0))
        
        if is_first and is_last:
            raise ValueError('Cannot specify both is_first=True and is_last=True')

        if is_last:
            ## Special case positioning the last milestone to be "right justified"
            left_offset_cm = self.ww_to_offset_cm(ww=milestone.ww_date) - width_cm + left_offset_cm
        else:
            left_offset_cm = self.ww_to_offset_cm(ww=milestone.ww_date) + left_offset_cm
        
        ### Update the milestone passed in with information related to rendering
        milestone.is_first = is_first
        milestone.is_last = is_last
        milestone.marker_left_cm = left_offset_cm
        milestone.marker_width_cm = width_cm

        add_rounded_rectangle_marker(shapes=self.slide.shapes, left_cm=left_offset_cm,\
            top_cm=vertical_offset_cm, width_cm = width_cm, height_cm=PRODUCT_MILESTONE_SHAPE_HEIGHT_CM,\
                text = shape_text, fill_rgb=milestone.fill_rgb, font_size=PRODUCT_MILESTONE_TEXT_IN_SHAPE_PTS,\
                     font_rgb_color=milestone.text_rgb, line_rgb=None)


        if milestone.adjacent_text != '':
        
            if is_last:
                ## Special case positioning the last milestone to be "right justified"
                text_width_cm = text_length_to_cm(milestone.adjacent_text, 8.0)
                left_offset_cm = left_offset_cm - text_width_cm - TWO_MM / 3.0
                text_align = PP_ALIGN.RIGHT
            else:
                text_width_cm = text_length_to_cm(milestone.text, 8.0)
                left_offset_cm = left_offset_cm + text_width_cm - TWO_MM / 2.0
                text_align = PP_ALIGN.LEFT
            
            ### Update the milestone passed in with information related to rendering
            milestone.text_box_left_cm = left_offset_cm
            milestone.text_box_width_cm = text_width_cm
            
            add_text_box(shapes=self.slide.shapes, left_cm=left_offset_cm, \
                top_cm=vertical_offset_cm, 
                width_cm = text_width_cm, \
                    height_cm=PRODUCT_MILESTONE_SHAPE_HEIGHT_CM,\
                    text = milestone.adjacent_text, text_align=text_align, text_auto_size=MSO_AUTO_SIZE.NONE,\
                    fill_rgb=None, line_rgb = None, line_width= Pt(1.0), line_dash = None,\
                    font_name='Intel Clear', font_size = Pt(8), font_rgb_color=RGBColor(0,0,0), font_bold=True, font_italic=False,\
                    font_underline=False, font_language_id=MSO_LANGUAGE_ID.NONE)

        #print(milestone, "width_cm",width_cm, "left_offset_cm", left_offset_cm)
        return milestone

    def render_annotations(self, pi, all_pis_dict):
        for key, ann in pi.annotations().items():
            ### Phase Implementation
            if isinstance(ann, an.Text):
                self._render_an_Text(pi=pi, ann=ann, all_pis_dict=all_pis_dict)
            
            if isinstance(ann, an.Delta):
                self._render_an_Delta(pi=pi, ann=ann, all_pis_dict=all_pis_dict)


    def _render_an_Text(self, pi, ann, all_pis_dict ):
        milestone_before = pi.milestone_before(col_name=ann.my_col)
        try:
            milestone_after = pi.milestone_after(col_name=ann.my_col)
        except:
            milestone_after = None

        shape_text = ann.text
        if milestone_after is not None:
            left_offset_cm , right_edge_cm = milestone_before.between_two_milestones(milestone_after)
        else:
            left_offset_cm  = milestone_before.right_of_marker_and_text()

        width_cm = 1.0

        vertical_offset_cm = pi.from_roadmap_top_cm + ADD_TO_CENTER_MILESTONE_IN_PRODUCT_SHAPE
        
        left_offset_cm = left_offset_cm + SM_CM  ## Fudge it a little

        add_text_box(shapes=self.slide.shapes, left_cm=left_offset_cm,\
                 top_cm=vertical_offset_cm, width_cm = width_cm, height_cm=PRODUCT_MILESTONE_SHAPE_HEIGHT_CM,\
                text = shape_text, text_align=PP_ALIGN.LEFT, text_auto_size=MSO_AUTO_SIZE.NONE,\
                fill_rgb=None, line_rgb = None, line_width= Pt(1.0), line_dash = None,\
                font_name='Intel Clear', font_size = Pt(8), font_rgb_color=RGBColor(0,0,0), font_bold=True)
    
    def _render_an_Delta(self, pi, ann, all_pis_dict ):
        milestone_before = pi.milestone_before(col_name=ann.my_col)
        milestone_after = pi.milestone_after(col_name=ann.my_col)
        
        #print(ann.my_col, milestone_before.text, milestone_before.ww_date, milestone_after.text, milestone_after.ww_date)
        
        ww_delta = milestone_before.ww_date.ww_delta_from(milestone_after.ww_date)

        if ann.date_format is None or ann.date_format == '' or \
            ann.date_format == an.DF_D:
            date_text = iw.wws_to_days_text(wws=ww_delta)
        elif ann.date_format == an.DF_WWs:
            date_text = iw.wws_to_text(wws=ww_delta)
        elif ann.date_format == an.DF_QTR:
            date_text = iw.wws_to_qtr_text(wws=ww_delta)
        elif ann.date_format == an.DF_YRS:
            date_text = iw.wws_to_years_text(wws=ww_delta)

        shape_text = ann.text + date_text + ann.after_text
        
        left_edge_cm, right_edge_cm = milestone_before.between_two_milestones(milestone_after)
        
        width_cm = 1.0
        
        left_offset_cm = left_edge_cm
        
        vertical_offset_cm = pi.from_roadmap_top_cm + ADD_TO_CENTER_MILESTONE_IN_PRODUCT_SHAPE*0.5
        
        add_text_box(shapes=self.slide.shapes, left_cm=left_offset_cm,\
                 top_cm=vertical_offset_cm, width_cm = width_cm, height_cm=PRODUCT_MILESTONE_SHAPE_HEIGHT_CM,\
                text = shape_text, text_align=PP_ALIGN.CENTER, text_auto_size=MSO_AUTO_SIZE.NONE,\
                fill_rgb=None, line_rgb = None, line_width= Pt(1.0), line_dash = None,\
                font_name='Intel Clear', font_size = Pt(9), font_rgb_color=RGBColor(0,0,0), font_bold=True )


    def annotate_slip(self, start_pi, start_ms, end_pi, end_ms, align_zero=False ):

        if align_zero:
            start_align_zero_ww = start_pi.first_milestone().ww_date
            end_align_zero_ww = end_pi.first_milestone().ww_date
        else:
            start_align_zero_ww = None
            end_align_zero_ww = None


        slip_wws = start_ms.ww_date.ww_delta_from(end_ms.ww_date)

        if slip_wws <= 0:
            print("Skiping annotation of slip of <= 0", slip_wws)
            return
        
        if slip_wws == 1:
            shape_text = f"{slip_wws}WW"
        else:
            shape_text = f"{slip_wws}WWs"

        left_offset_cm = self.ww_to_offset_cm(ww=start_ms.ww_date, \
            align_zero_ww=start_align_zero_ww)
        
        right_offset_cm = self.ww_to_offset_cm(ww=end_ms.ww_date, \
            align_zero_ww=end_align_zero_ww)

        width_cm = right_offset_cm - left_offset_cm

        vertical_offset_cm = start_pi.from_roadmap_top_cm

        add_rounded_rectangle_marker(shapes=self.slide.shapes, left_cm=left_offset_cm,\
            top_cm=vertical_offset_cm, width_cm = width_cm, height_cm=PRODUCT_MILESTONE_SHAPE_HEIGHT_CM,\
                mso_shape=MSO_SHAPE.PENTAGON,
                text = shape_text, fill_rgb=RGBColor(128,0,0), font_size=PRODUCT_MILESTONE_TEXT_IN_SHAPE_PTS,\
                     font_rgb_color=RGBColor(255,255,255), line_rgb=None)
        
        add_connector(shapes=self.slide.shapes, x1=right_offset_cm, \
            y1=vertical_offset_cm+PRODUCT_SHAPE_HEIGHT_CM/2.0,\
            x2=right_offset_cm, y2=end_pi.from_roadmap_top_cm+PRODUCT_SHAPE_HEIGHT_CM/2.0)

        return
    
    def annotate_inline(self, pi, align_zero=False ):
        vertical_offset_cm = pi.from_roadmap_top_cm

        print('Annotating:',pi.shorthand_name)
        if align_zero:
            align_zero_ww = pi.first_milestone().ww_date
        else:
            align_zero_ww = None
        
        for i, ms in enumerate(pi.milestones):
            if i == len(pi.milestones)-1:
                break
            ms_next = pi.milestones[i+1]

            ww_delta = ms.ww_date.ww_delta_from(ms_next.ww_date)
            shape_text = str(ww_delta)
            left_edge_cm, right_edge_cm = ms.between_two_milestones(ms_next)
            width_cm = 1.0
            left_offset_cm = left_edge_cm
            
            add_text_box(shapes=self.slide.shapes, left_cm=left_offset_cm,\
                 top_cm=vertical_offset_cm, width_cm = width_cm, height_cm=PRODUCT_MILESTONE_SHAPE_HEIGHT_CM,\
                text = shape_text, text_align=PP_ALIGN.CENTER, text_auto_size=MSO_AUTO_SIZE.NONE,\
                fill_rgb=None, line_rgb = None, line_width= Pt(1.0), line_dash = None,\
                font_name='Intel Clear', font_size = Pt(12), font_rgb_color=RGBColor(0,0,0) )
    
        return

def add_rounded_rectangle_marker(shapes, left_cm, top_cm, width_cm, height_cm=PRODUCT_SHAPE_HEIGHT_CM,\
         mso_shape=MSO_SHAPE.ROUNDED_RECTANGLE,\
         text = None, text_align=PP_ALIGN.CENTER, text_auto_size=MSO_AUTO_SIZE.NONE,\
        fill_rgb=None, line_rgb = None, line_width= Pt(1.0), line_dash = None,\
             font_name='Arial', font_size = Pt(12), font_rgb_color=RGBColor(0,0,0), font_bold=False, font_italic=False,\
             font_underline=False, font_language_id=MSO_LANGUAGE_ID.NONE, adjustments = [1.0] ):
    """
    Function for drawing product indicators on roadamp pptxs
    """
    shape = shapes.add_shape(mso_shape, Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm))
    
    ### Manage adjustments (fully round the end...)
    for i, adj in enumerate(adjustments):
        shape.adjustments[i] = adj

    ### Manage the FillFormat properties
    fill = shape.fill
    if fill_rgb is not None:
        fill.solid()
        fill.fore_color.rgb = fill_rgb
    else:
        fill.background()  # For transparent shapes - just here for reference

    ### Manage the LineFormat properties
    line = shape.line

    """
    MSO_LINE  for DASH STYLE
    DASH - Line consists of dashes only.
    DASH_DOT - Line is a dash-dot pattern.
    DASH_DOT_DOT - Line is a dash-dot-dot pattern.
    LONG_DASH - Line consists of long dashes.
    LONG_DASH_DOT -Line is a long dash-dot pattern.
    ROUND_DOT - Line is made up of round dots.
    SOLID - Line is solid.
    SQUARE_DOT - Line is made up of square dots.
    DASH_STYLE_MIXED - Not supported.
    """
    if line_rgb is not None:
        line.color.rgb = line_rgb
        line.width = line_width
        line.dash_style = line_dash
    else:
        line.fill.background()  # For transparent shapes - just here for reference
    
    """
    shape.text_frame TextFrame Object
    """
    # TextFrame
    if text is not None:
        text_frame = shape.text_frame
        text_frame.clear()    # Premptively remove any paragraphs

        """
        auto_size: MSO_AUTO_SIZE.NONE, MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT, or MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE.
        """
        text_frame.auto_size = text_auto_size

        text_frame.margin_left = Cm(0.0)
        text_frame.margin_right = Cm(0.0)
        text_frame.margin_bottom = Cm(0.0)
        text_frame.margin_top = Cm(0.0)
        
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = text_align
        set_font(font=paragraph.font, name=font_name, size = font_size, rgb_color=font_rgb_color,\
            bold=font_bold, italic=font_italic, underline=font_underline, language_id=font_language_id )
        paragraph.text = text

    return shape

def set_font(font, name, size = Pt(12), rgb_color=RGBColor(0,0,0), bold=False,\
     italic=False, underline=False, language_id=MSO_LANGUAGE_ID.NONE ):
    font.name = name
    font.size = size
    font.color.rgb = rgb_color
    font.bold = bold
    font.italic = italic
    font.language_id = language_id
    font.underline = underline
    return font

def add_text_box(shapes, left_cm, top_cm, width_cm, height_cm=PRODUCT_SHAPE_HEIGHT_CM,\
         text = None, text_align=PP_ALIGN.CENTER, text_auto_size=MSO_AUTO_SIZE.NONE,\
        fill_rgb=None, line_rgb = None, line_width= Pt(1.0), line_dash = None,\
             font_name='Arial', font_size = Pt(12), font_rgb_color=RGBColor(0,0,0), font_bold=False, font_italic=False,\
             font_underline=False, font_language_id=MSO_LANGUAGE_ID.NONE, adjustments = [1.0] ):
    """
    Function for drawing product indicators on roadamp pptxs
    """
    shape = shapes.add_textbox(Cm(left_cm), Cm(top_cm),\
         Cm(width_cm), Cm(height_cm))

    ### Manage the FillFormat properties
    fill = shape.fill
    if fill_rgb is not None:
        fill.solid()
        fill.fore_color.rgb = fill_rgb
    else:
        fill.background()  # For transparent shapes - just here for reference

    ### Manage the LineFormat properties
    line = shape.line

    """
    MSO_LINE  for DASH STYLE
    DASH - Line consists of dashes only.
    DASH_DOT - Line is a dash-dot pattern.
    DASH_DOT_DOT - Line is a dash-dot-dot pattern.
    LONG_DASH - Line consists of long dashes.
    LONG_DASH_DOT -Line is a long dash-dot pattern.
    ROUND_DOT - Line is made up of round dots.
    SOLID - Line is solid.
    SQUARE_DOT - Line is made up of square dots.
    DASH_STYLE_MIXED - Not supported.
    """
    if line_rgb is not None:
        line.color.rgb = line_rgb
        line.width = line_width
        line.dash_style = line_dash
    else:
        line.fill.background()  # For transparent shapes - just here for reference
    
    """
    shape.text_frame TextFrame Object
    """
    # TextFrame
    if text is not None:
        text_frame = shape.text_frame
        text_frame.clear()    # Premptively remove any paragraphs

        """
        auto_size: MSO_AUTO_SIZE.NONE, MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT, or MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE.
        """
        text_frame.auto_size = text_auto_size

        text_frame.margin_left = Cm(0.0)
        text_frame.margin_right = Cm(0.0)
        text_frame.margin_bottom = Cm(0.0)
        text_frame.margin_top = Cm(0.0)
        
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = text_align
        set_font(font=paragraph.font, name=font_name, size = font_size, rgb_color=font_rgb_color,\
            bold=font_bold, italic=font_italic, underline=font_underline, language_id=font_language_id )
        paragraph.text = text

    return shape

def add_connector(shapes, x1, y1, x2, y2, shadow=False ):
    
    connector = shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(x1), Cm(y1),Cm(x2), Cm(y2))
    color_line = connector.line
    color_line.width=Cm(0.02)
    color_line.fill.solid()
    color_line.fill.fore_color.rgb=RGBColor(0,0,0)
    color_line.dash_style=MSO_LINE.DASH

    #connector.line.fill.fore_color.rgb=RGBColor(0,0,0)

    #connector.line.fill.fore_color= RGBColor(0,0,0)
    
    #if shadow is False:
    #    connector.shadow.inherit=False
    
    return connector  
"""
class RoadmapPPT:

    def __init__(self, path_to_template, slide_config_list):
        self.path_to_template = path_to_template
        self.g_ppt = Presentation(path_to_template)
        self.slide_height=self.g_ppt.slide_height
        self.slide_width = self.g_ppt.slide_width
        self.slides = self.g_ppt.slides

        self.layouts = {}
        for layout in self.g_ppt.slide_layouts:
            print(layout.name)
            self.layouts[layout.name]=layout
    
    def add_slide(self, layout):
        new_slide = self.slides.add_slide(self.layouts[layout])
        return new_slide
    
    def dump(self):
        print(self.slide_height, self.slide_width)
        return
    
    def save(self, full_output_path ):
        self.g_ppt.slide_height=self.slide_height
        self.g_ppt.slide_width=self.slide_width
        self.g_ppt.save(full_output_path)
        return
    

    def draw_shape(self, slide ):
        'ACTION_BUTTON_BACK_OR_PREVIOUS', 'ACTION_BUTTON_BEGINNING', 'ACTION_BUTTON_CUSTOM', 'ACTION_BUTTON_DOCUMENT',
        'ACTION_BUTTON_END', 'ACTION_BUTTON_FORWARD_OR_NEXT', 'ACTION_BUTTON_HELP', 'ACTION_BUTTON_HOME', 
        'ACTION_BUTTON_INFORMATION', 'ACTION_BUTTON_MOVIE', 'ACTION_BUTTON_RETURN', 'ACTION_BUTTON_SOUND', 
        'ARC', 'BALLOON', 'BENT_ARROW', 'BENT_UP_ARROW', 'BEVEL', 'BLOCK_ARC', 'CAN', 'CHART_PLUS', 
        'CHART_STAR', 'CHART_X', 'CHEVRON', 'CHORD', 'CIRCULAR_ARROW', 'CLOUD', 'CLOUD_CALLOUT', 
        'CORNER', 'CORNER_TABS', 'CROSS', 'CUBE', 'CURVED_DOWN_ARROW', 'CURVED_DOWN_RIBBON', 
        'CURVED_LEFT_ARROW', 'CURVED_RIGHT_ARROW', 'CURVED_UP_ARROW', 'CURVED_UP_RIBBON', 'DECAGON',
        'DIAGONAL_STRIPE', 'DIAMOND', 'DODECAGON', 'DONUT', 'DOUBLE_BRACE', 'DOUBLE_BRACKET', 'DOUBLE_WAVE',
        'DOWN_ARROW', 'DOWN_ARROW_CALLOUT', 'DOWN_RIBBON', 'EXPLOSION1', 'EXPLOSION2', 
        'FLOWCHART_ALTERNATE_PROCESS', 'FLOWCHART_CARD', 'FLOWCHART_COLLATE', 'FLOWCHART_CONNECTOR', 
        'FLOWCHART_DATA', 'FLOWCHART_DECISION', 'FLOWCHART_DELAY', 'FLOWCHART_DIRECT_ACCESS_STORAGE',
            'FLOWCHART_DISPLAY', 'FLOWCHART_DOCUMENT', 'FLOWCHART_EXTRACT', 'FLOWCHART_INTERNAL_STORAGE', 
            'FLOWCHART_MAGNETIC_DISK', 'FLOWCHART_MANUAL_INPUT', 'FLOWCHART_MANUAL_OPERATION', 'FLOWCHART_MERGE', 
            'FLOWCHART_MULTIDOCUMENT', 'FLOWCHART_OFFLINE_STORAGE', 'FLOWCHART_OFFPAGE_CONNECTOR', 'FLOWCHART_OR', 
            'FLOWCHART_PREDEFINED_PROCESS', 'FLOWCHART_PREPARATION', 'FLOWCHART_PROCESS', 'FLOWCHART_PUNCHED_TAPE',
            'FLOWCHART_SEQUENTIAL_ACCESS_STORAGE', 'FLOWCHART_SORT', 'FLOWCHART_STORED_DATA', 'FLOWCHART_SUMMING_JUNCTION', 
            'FLOWCHART_TERMINATOR', 'FOLDED_CORNER', 'FRAME', 'FUNNEL', 'GEAR_6', 'GEAR_9', 'HALF_FRAME', 'HEART', 'HEPTAGON', 
            'HEXAGON', 'HORIZONTAL_SCROLL', 'ISOSCELES_TRIANGLE', 'LEFT_ARROW', 'LEFT_ARROW_CALLOUT', 'LEFT_BRACE', 'LEFT_BRACKET',
            'LEFT_CIRCULAR_ARROW', 'LEFT_RIGHT_ARROW', 'LEFT_RIGHT_ARROW_CALLOUT', 'LEFT_RIGHT_CIRCULAR_ARROW', 
            'LEFT_RIGHT_RIBBON', 'LEFT_RIGHT_UP_ARROW', 'LEFT_UP_ARROW', 'LIGHTNING_BOLT', 'LINE_CALLOUT_1', 
            'LINE_CALLOUT_1_ACCENT_BAR', 'LINE_CALLOUT_1_BORDER_AND_ACCENT_BAR', 'LINE_CALLOUT_1_NO_BORDER', 
            'LINE_CALLOUT_2', 'LINE_CALLOUT_2_ACCENT_BAR', 'LINE_CALLOUT_2_BORDER_AND_ACCENT_BAR', 'LINE_CALLOUT_2_NO_BORDER', 
            'LINE_CALLOUT_3', 'LINE_CALLOUT_3_ACCENT_BAR', 'LINE_CALLOUT_3_BORDER_AND_ACCENT_BAR', 'LINE_CALLOUT_3_NO_BORDER', 
            'LINE_CALLOUT_4', 'LINE_CALLOUT_4_ACCENT_BAR', 'LINE_CALLOUT_4_BORDER_AND_ACCENT_BAR', 'LINE_CALLOUT_4_NO_BORDER', 
            'LINE_INVERSE', 'MATH_DIVIDE', 'MATH_EQUAL', 'MATH_MINUS', 'MATH_MULTIPLY', 'MATH_NOT_EQUAL', 'MATH_PLUS', 'MOON', 
            'NON_ISOSCELES_TRAPEZOID', 'NOTCHED_RIGHT_ARROW', 'NO_SYMBOL', 'OCTAGON', 'OVAL', 'OVAL_CALLOUT', 'PARALLELOGRAM', 
            'PENTAGON', 'PIE', 'PIE_WEDGE', 'PLAQUE', 'PLAQUE_TABS', 'QUAD_ARROW', 'QUAD_ARROW_CALLOUT', 'RECTANGLE', 
            'RECTANGULAR_CALLOUT', 'REGULAR_PENTAGON', 'RIGHT_ARROW', 'RIGHT_ARROW_CALLOUT', 'RIGHT_BRACE', 'RIGHT_BRACKET', 
            'RIGHT_TRIANGLE', 'ROUNDED_RECTANGLE', 'ROUNDED_RECTANGULAR_CALLOUT', 'ROUND_1_RECTANGLE', 'ROUND_2_DIAG_RECTANGLE', 
            'ROUND_2_SAME_RECTANGLE', 'SMILEY_FACE', 'SNIP_1_RECTANGLE', 'SNIP_2_DIAG_RECTANGLE', 'SNIP_2_SAME_RECTANGLE', 
            'SNIP_ROUND_RECTANGLE', 'SQUARE_TABS', 'STAR_10_POINT', 'STAR_12_POINT', 'STAR_16_POINT', 'STAR_24_POINT', 
            'STAR_32_POINT', 'STAR_4_POINT', 'STAR_5_POINT', 'STAR_6_POINT', 'STAR_7_POINT', 'STAR_8_POINT', 'STRIPED_RIGHT_ARROW', 
            'SUN', 'SWOOSH_ARROW', 'TEAR', 'TRAPEZOID', 'UP_ARROW', 'UP_ARROW_CALLOUT', 'UP_DOWN_ARROW', 'UP_DOWN_ARROW_CALLOUT', 
            'UP_RIBBON', 'U_TURN_ARROW', 'VERTICAL_SCROLL', 'WAVE', 
        shapes = slide.shapes
        
        left = Inches(1.0)
        right = Inches(1.0)
        width = Inches(1.0)
        height = Inches(1.0)
        
        shape=shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left,right,width, height)

        fill = False
        shadow = False
        
        if fill is False:
            shape.fill.background()
        
        if shadow is False:
            shape.shadow.inherit=False
        
        return shape
    
def draw_marker(x, y, d, shapes, fill=False, shadow=False ):
    left = Inches(x)
    right = Inches(y)
    width = Inches(d)
    height = Inches(d)
    
    shape=shapes.add_shape(MSO_SHAPE.OVAL, left,right,width, height)
    
    if fill is False:
        shape.fill.background()
    
    if shadow is False:
        shape.shadow.inherit=False
    
    return shape

def draw_text(x, y, text, shapes, size=12, shadow=False ):
    left = Inches(x)
    right = Inches(y)
    width = Inches(0.1)
    height = Inches(0.1)
    
    text_shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, right ,width, height)
    text_shape.fill.background()
    text_shape.line.fill.background()
    
    if shadow is False:
        text_shape.shadow.inherit=False

    text_frame = text_shape.text_frame

    text_frame.clear()  # not necessary for newly-created shape
    text_frame.auto_size=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    text_frame.word_wrap=False

    p = text_frame.paragraphs[0]
    p.alignment=PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text

    font = run.font
    font.name = 'Calibri'
    font.size = Pt(size)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    font.color.theme_color = MSO_THEME_COLOR.TEXT_1
    
    return text_shape

 def draw_connector(x1, y1, x2, y2, shapes, shadow=False ):
    
    connector = shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1),Inches(x2), Inches(y2))
    
    connector.line.width=Inches(0.01)
    #connector.fill.background()
    #connector.connector.line.color= MSO_THEME_COLOR.TEXT_1
    
    if shadow is False:
        connector.shadow.inherit=False
    
    return connector  

"""
if __name__ == '__main__':
    ml_list = []

    for i in range(10):
        ms = Milestone(text=str(i), ww_date = iw.random_WW(min_year=19, max_year=24))
        ml_list.append(ms)
    
    print([str(m) for m in ml_list])
    ml_list = sorted(ml_list)
    print([str(m) for m in ml_list]) 

    """
    IAGS_PPTX = '/Users/scotttan/Intel Corporation/Graphics Golden Doc Development - RoadmapPPTGeneration/ScottsWorkbench/BlankIAGSTemplate.pptx'
    IAGS_PPTX_OUT = '/Users/scotttan/Intel Corporation/Graphics Golden Doc Development - RoadmapPPTGeneration/ScottsWorkbench/IAGS_OUT.pptx'
    
    p = Presentation(IAGS_PPTX)
    slide = p.slides[0]

    start_ww = iw.WW(ww=13, year=20)
    end_ww = iw.WW(ww=13, year=24)
    grid = CanvasGrid(start_ww=start_ww, end_ww=end_ww)
    print(grid.total_quarters())
    
    rc = RoadmapCanvas(grid=grid)
    rc.render_yrs_qts_table( shapes=slide.shapes )
    top_of_draw_area = rc.top_of_draw_area
    height_of_draw_area = rc.draw_area_height

    for i, ww_offset in enumerate(range(rc.grid.ww_span)):
        new_ww = rc.grid.left_edge_ww.add_wws_return_new(ww_offset)
        x_cm = rc.ww_to_slide_x_cm(ww=new_ww)
        #print(i, new_ww, f"{x_cm:.2f}" )
        rc.test_circle(shapes=slide.shapes,abs_x_cm=x_cm,text=str(new_ww.ww))

    h = OrderedDict()
    h['Data Center']=OrderedDict()
    h['Data Center']['Compute']='Compute'
    h['Data Center']['ComputePH']=''
    h['Data Center']['MultiPurpose']='Multi Purpose'
    h['Discrete Client']=OrderedDict()
    h['Discrete Client']['Pro']='Pro'
    h['Discrete Client']['Enthusiast']='Enthusiast'
    h['Discrete Client']['Performance']='Performance'
    h['Discrete Client']['Mainstream']='Mainstream'
    h['Discrete Client']['Entry']='Entry'

    slg = SwimlanesGrid(hierarchy=h)
    slt = SwimlaneTable(grid=slg,table_top_cm=top_of_draw_area,\
        table_left_edge_cm=0.0, table_height_cm=height_of_draw_area)
    
    for biz in slg.businesses():
        for seg in slg.segments(biz):
            abs_y_cm = slt.y_value_in_cm_for_segment(business=biz, segment=seg)
            print(biz, seg, abs_y_cm)
            rc.test_circle(shapes=slide.shapes,abs_x_cm=3.5,abs_y_cm=abs_y_cm,\
                text=seg, width_cm=3.0,font_size=Pt(12))

    
    slt.render_swimlanes_table(shapes=slide.shapes)
    
    p.save(IAGS_PPTX_OUT)
    """

    """
    import num2words
    GOLDEN_TABLE_PPTX = '/Users/scotttan/Intel Corporation/Graphics Golden Doc Development - RoadmapPPTGeneration/ScottsWorkbench/JustQuarters.pptx'
    GOLDEN_TABLE_PPTX_OUT = '/Users/scotttan/Intel Corporation/Graphics Golden Doc Development - RoadmapPPTGeneration/ScottsWorkbench/JustQuarters_OUT.pptx'

    milestones = []
    start_ww= iw.WW(ww=1,year=20)
    step_ww = 20

    for i in range(8):
        start_ww.add_wws(step_ww)
        milestones.append(Milestone(text=str(start_ww), ww_date=start_ww.copy()))

    for ms in milestones:
        print(ms)
    
    pi = ProgramInformation(name='FirstProgram', milestones_list = milestones)
    print(pi, pi.first_milestone(), pi.last_milestone())

    p = Presentation(GOLDEN_TABLE_PPTX)

    golden_slide = p.slides[0]

    rs = RoadmapSlide(slide=golden_slide, roadmap_left_edge_cm=ROADMAP_ZERO_LEFT, \
                     roadmap_top_cm=ROADMAP_TOP, roadmap_width_cm=ROADMAP_WIDTH,\
                    name_left_edge_cm=NAME_LEFT_EDGE,year_zero = 19, quarter_span=20)
    print(rs)

    rs.render_program(pi)

    p.save(GOLDEN_TABLE_PPTX_OUT)
    """
    """
    start_ww = iw.WW(ww=13, year=19)
    end_ww = iw.WW(ww=14, year=21)
    grid = YrQtWWGrid(start_ww=start_ww, end_ww=end_ww)
    print(grid,'\n', grid.quarters_in_first_year(),  grid.total_quarters())
    """











