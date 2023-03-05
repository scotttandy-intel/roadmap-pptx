"""
roadmap_json.py
Deconstruct all (required) pptx objects into a JSON representation

Author Scott Tandy
Copyright Intel Corporation 2020
"""

import pptx


def layouts_to_dict( slide_layouts):
    layout_dict = {}
    for layout in slide_layouts:
        layout_dict[layout.name]=layout
    return layout_dict


def show_slides( slides ):
    for i, slide in enumerate(slides):
        print(i, slide.slide_layout.name)


def walk_object(obj, ignore_underscore=True, ignore_methods=True, depth=0, to_depth=2, skip_list=None, no_walk=None):
    for attr in dir(obj):
        if len(attr) == 0:
            continue  ##Protects indexing;

        if ignore_underscore and attr[0]=='_':
            continue
        
        try:
            attr_obj = getattr(obj, attr)
        except:
            print('-'*depth + attr + ':' + 'ERROR GETTING ATTRIBUTE')
            continue

        if ignore_methods and callable(attr_obj):
            continue

        if skip_list is not None:
             if type(attr_obj) in skip_list:
                 continue
        
        if isinstance(attr_obj, str):
            str_value = attr_obj
        elif isinstance(attr_obj, int):
            str_value = str(attr_obj)
        elif isinstance(attr_obj, float):
            str_value = str(attr_obj)
        else:
            str_value = str(attr_obj)

        print('-'*depth + attr + ':' + str(type(attr_obj)) + ':' + str_value)
        
        if no_walk is not None and type(attr_obj) in no_walk:
            continue

        if depth < to_depth:
            if isinstance(attr_obj, list):
                print('-'*depth +'[list]')
                for list_obj in attr_obj:
                    print('-'*depth +'[listitem]')
                    walk_object(obj=list_obj, depth=depth+1, to_depth=to_depth, skip_list=skip_list, no_walk=no_walk)
            elif isinstance(attr_obj, pptx.shapes.shapetree.SlideShapes):
                print('-'*depth +'[slide shape list]')
                for list_obj in attr_obj:
                    print('-'*depth +'[shape]')
                    walk_object(obj=list_obj, depth=depth+1, to_depth=to_depth, skip_list=skip_list, no_walk=no_walk)
            else:
                walk_object(obj=attr_obj, depth=depth+1, to_depth=to_depth, skip_list=skip_list, no_walk=no_walk)
        

def dump_table(table_obj):
    for col in table_obj.columns:
        print("{:.2f}".format(col.width/pptx.util.Inches(1.0)) + ' ', end='')
    print('')
    
    for row in table_obj.rows:
        print("Row Height Inches{:.2f}".format(row.height / pptx.util.Inches(1.0)))
        for cell in row.cells:
            print(cell.text, '-',end='')
        print('')
        for cell in row.cells:
            print(str(cell.is_merge_origin) + '-',end='')
        print('')

def print_obj(obj,depth=0):
    if obj is None:
        print("None")
    else:
        print(obj)

def print_header(text, depth=0):
    print('\n'+str(depth)+'_'*depth+'------------------------------')
    print(str(depth)+'_'*depth+text)
    print(str(depth)+'_'*depth+'------------------------------')
"""
Presentation
"""
def presentation_from_pptx( path_to_pptx):
    return pptx.Presentation(pptx=path_to_pptx)


def dump_pptx_object(obj, serialize, depth=0):
    print(str(depth)+'_'*depth+"OBJECT: ", obj)
    
    for attr in dir(obj):
        if attr not in serialize:
            #print("Skipping: ", attr)
            continue
        try:
            attr_obj = getattr(obj, attr)
        except:
            print(attr + "**Error getting attribute")
            continue
            
        if serialize[attr] is not None:
            try:
                print(attr,end=':')
                serialize[attr](attr_obj)
            except:
                print("Error calling output for attribute:", attr)
    return

def dump_presentation( presentation):
    serialize = {
        'core_properties':dump_core_properties,
        'element':None,
        'notes_master':None,
        'part':None,
        'save':None,
        'slide_height':print_obj,
        'slide_layouts':dump_slide_layouts,
        'slide_master':None,
        'slide_masters':None,
        'slide_width':print_obj,
        'slides':dump_slides,
        }
    print_header('Presentation')
    dump_pptx_object(obj=presentation, serialize=serialize)
    
    return

def dump_core_properties(core_properties):
    serialize = {\
    'after_unmarshal':None,
    'author':print_obj,
    'before_marshal':None,
    'blob':None,
    'category':print_obj,
    'comments':print_obj,
    'content_status':print_obj,
    'content_type':print_obj,
    'created':print_obj,
    'default':None,
    'drop_rel':None,
    'identifier':print_obj,
    'keywords':print_obj,
    'language':print_obj,
    'last_modified_by':print_obj,
    'last_printed':print_obj,
    'load':None,
    'load_rel':None,
    'modified':print_obj,
    'package':None,
    'part':None,
    'part_related_by':None,
    'partname':print_obj,
    'relate_to':None,
    'related_parts':print_obj,
    'rels':print_obj,
    'revision':print_obj,
    'subject':print_obj,
    'target_ref':None,
    'title':print_obj,
    'version':print_obj,
    }
    print_header('Core Properties')
    dump_pptx_object(obj=core_properties, serialize=serialize)
    return

def dump_slide_layouts(slide_layouts):
    print_header('Slide Layouts')
    for layout in slide_layouts:
        print(layout.name)
    return

def dump_slides(slides):
    print_header('Slides')
    for slide in slides:
        dump_slide(slide)


    #########   REMOVE TO DUMP ALL SLIDES!
        break
    ##########
    return

def dump_slide(slide):
    serialize = {
        'background':None,
        'element':None,
        'follow_master_background':print_obj,
        'has_notes_slide':print_obj,
        'name':print_obj,
        'notes_slide':None,
        'part':print_obj,
        'placeholders':dump_slide_placeholders,
        'shapes':dump_slide_shapes,
        'slide_id':print_obj,
        'slide_layout':print_obj
    }
    print_header('Slide')
    dump_pptx_object( slide, serialize=serialize)
    return

def dump_slide_placeholders(placeholders):
    print_header('Slide Placeholders')
    return

def dump_slide_shapes( shapes, depth=0 ):
    print_header('Slide Shapes',depth=depth)

    for shape in shapes:
        if 'shapes' in dir(shape):
            dump_slide_shapes(shape.shapes, depth=depth+1)
        
        dump_slide_shape(shape, depth=depth)

def dump_slide_shape(shape, depth=0):
    serialize = {
        pptx.shapes.shapetree.GroupShape:dump_slide_type_group_shape,
        pptx.shapes.autoshape.Shape:dump_slide_type_shape,
        pptx.shapes.connector.Connector:print_obj,
        pptx.shapes.freeform.FreeformBuilder:print_obj,
        pptx.shapes.picture.Picture:print_obj,
        pptx.shapes.graphfrm.GraphicFrame:print_obj
    }

    if type(shape) in serialize:
        serialize[type(shape)](shape, depth=depth)

    
    return

def dump_slide_type_group_shape(groupshape, depth=0):
    print_header('Group Shape')
    serialize = \
        {
        'has_chart':print_obj,
        'has_table':print_obj,
        'height':print_obj,
        'is_placeholder':print_obj,
        'left':print_obj,
        'name':print_obj,
        'rotation':print_obj,
        'shadow':dump_shadow_format,
        'shape_id':print_obj,
        'shape_type':print_obj,
        'shapes':print_obj,
        'top':print_obj,
        'width':print_obj
    }

    dump_pptx_object(groupshape, serialize, depth=depth)
    return

def dump_shadow_format(shadow_format, depth=0):
    print_header('Shadow Format')
    serialize = \
        {
             'inherit':print_obj
        }
        
    dump_pptx_object(shadow_format, serialize, depth=depth)
    return

def dump_slide_type_shape(groupshape, depth=0):
    print_header('Shape')
    serialize = \
        {
        'adjustments':dump_slide_type_shape_adjustments,
        'auto_shape_type':print_obj,
        'click_action':dump_slide_type_shape_action_setting,
        'fill':dump_slide_type_shape_fill_format,
        'has_chart':print_obj,
        'has_table':print_obj,
        'has_text_frame':print_obj,
        'height':print_obj,
        'is_placeholder':print_obj,
        'left':print_obj,
        'line':dump_line_format,
        'ln':print_obj,
        'name':print_obj,
        'part':print_obj,
        'placeholder_format':print_obj,
        'rotation':print_obj,
        'shadow':print_obj,
        'shape_id':print_obj,
        'shape_type':print_obj,
        'text':print_obj,
        'text_frame':print_obj,
        'top':print_obj,
        'width':print_obj
    }

    dump_pptx_object(groupshape, serialize, depth=depth)

    return

def dump_slide_type_shape_adjustments(adjustments, depth=0):
    print_header("Adjustments", depth)
    for adjustment in adjustments:
        print_obj(adjustment)

def dump_slide_type_shape_action_setting(action_setting, depth=0):
    print_header("Action Setting", depth)
    serialize = \
        {
             'inherit':print_obj
        }
        
    dump_pptx_object(action_setting, serialize, depth=depth)
    
    return

def dump_slide_type_shape_fill_format(fill_format, depth=0):
    print_header("Fill Format", depth)
    serialize = \
        {
        'back_color':dump_color_format,
        'fore_color':dump_color_format,
        'gradient_angle':print_obj,
        'gradient_stops':print_obj,
        'pattern':print_obj,
        'type':print_obj
        }
        
    dump_pptx_object(fill_format, serialize, depth=depth)
    
    return

def dump_color_format( color_format, depth=0):
    print_header("Color Format", depth)
    serialize = \
        {
        'brightness':print_obj,
        'rgb':print_obj,
        'theme_color':print_obj,
        'type':print_obj
        }
    dump_pptx_object(color_format, serialize, depth=depth)

def dump_line_format( line_format, depth=0):
    print_header("Line Format", depth)
    serialize = \
        {
        'color':dump_color_format,
        'dash_style':print_obj,
        'fill':dump_slide_type_shape_fill_format,
        'width':print_obj
        }
   
    dump_pptx_object(line_format, serialize, depth=depth)

if __name__ == '__main__':
    

    



