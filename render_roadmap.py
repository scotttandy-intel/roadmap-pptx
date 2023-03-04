"""
render_roadmap.py

Collection of functions to combine the roadmap_table.py (and helper) information and 
render differently formated roadmaps onto the roadmap_pptx.py presentations
"""

###
from collections import OrderedDict
from pptx import Presentation

import intel_roadmap as ir
import intel_ww as iw
import roadmap_config as rc
import roadmap_table as rt
import roadmap_helper as rh
import roadmap_pptx as rp
import annotations as an


NOT_PROVIDED = "???"
PROG_PROG_VERT_DISTANCE_CM = 0.5

def render_swimlanes_and_qts(roadmap_slide, roadmap_canvas, swimlane_table ):
    roadmap_canvas.render_yrs_qts_table(shapes=roadmap_slide.slide.shapes)
    swimlane_table.render_swimlanes_table(shapes=roadmap_slide.slide.shapes)

def render_roadmap( roadmap_slide, roadmap_configuration, roadmap_table, \
    swimlane_table, align_zero=False ):
    
    """
    Returns simple_name_dict
    """
    if not roadmap_table.is_concept_format:

        ### Holds Program Informations based on simple name - useful to know where they were rendered
        shorthand_name_dict = {}   
        
        for biz in roadmap_table.business_list():
            for ss in roadmap_table.simple_segment_list(business=biz):
                rows = roadmap_table.by_type_simplesegment(business=biz, simple_segment=ss, type_name=rt.SI_PRODUCT )
                
                if rows.shape[0] == 0:
                    continue
                
                row_count = 0
                
                for i, row in rows.iterrows():
                    try:
                        y_offset = swimlane_table.y_value_in_cm_for_segment(business=biz, segment=ss)
                    except:
                        print('Skipping:',biz, ss)
                        continue
                    
                    pi = rp.rt_to_ProgramInformation(roadmap_table=roadmap_table, roadmap_configuration=roadmap_configuration,\
                        shorthand_name=row[rt.SHORTHAND_NAME], from_roadmap_top_cm=y_offset + (row_count * PROG_PROG_VERT_DISTANCE_CM),\
                            missing_ms=NOT_PROVIDED)
                    
                    ### Note - the render program function actually updates each of the
                    ### milestones with x axis information about where they were rendered.
                    pi = roadmap_slide.render_program(program_information=pi, align_zero=align_zero)               
                    
                    shorthand_name_dict[pi.shorthand_name] = pi

                    row_count = row_count + 1
        
    else:
        shorthand_name_dict = render_roadmap_concept(roadmap_slide=roadmap_slide, roadmap_configuration=roadmap_configuration,\
            roadmap_table=roadmap_table, swimlane_table=swimlane_table, align_zero=align_zero)
    
    return shorthand_name_dict

def add_to_ordered_dict(d, key1, key2, key3, obj):
    if key1 not in d.keys():
        d[key1] = OrderedDict()
    if key2 not in d[key1].keys():
        d[key1][key2]=OrderedDict()
    
    d[key1][key2][key3] = obj
    
    return d

def render_roadmap_concept( roadmap_slide, roadmap_configuration, roadmap_table, \

    swimlane_table, align_zero=False ):

   #major_col_in_df_name = roadmap_configuration.roadmap_columns[roadmap_configuration.major_column_name]
    #minor_col_in_df_name = roadmap_configuration.roadmap_columns[roadmap_configuration.minor_column_name]
    #product_name_col_in_df = roadmap_configuration.roadmap_columns[rt.CT_NAME]

    major_col_in_df_name = roadmap_configuration.major_column_name
    minor_col_in_df_name = roadmap_configuration.minor_column_name
    product_name_col_in_df = roadmap_configuration.name_col
            
    shorthand_name_dict = {}   
    
    for major_key in roadmap_configuration.swimlanes_hierarchy.keys():
        print(major_key)
        for minor_key in roadmap_configuration.swimlanes_hierarchy[major_key]:
            print('*'*2, minor_key, roadmap_configuration.swimlanes_hierarchy[major_key][minor_key])

            major_minor_tuples = [ (major_col_in_df_name,[major_key]),\
                (minor_col_in_df_name,[minor_key]) ]
            
            rows = roadmap_table.rows_isin_col_tuples(col_tuples=major_minor_tuples)

            if rows.shape[0] == 0:
                    continue
            
            row_count = 0
                
            for i, row in rows.iterrows():
                print(row[major_col_in_df_name], row[minor_col_in_df_name], \
                    row[product_name_col_in_df])
                try:
                    y_offset = \
                        swimlane_table.y_value_in_cm_for_major_column_minor_column(major_value=row[major_col_in_df_name], \
                            minor_value=row[minor_col_in_df_name])
                except:
                    print('Skipping:',row[major_col_in_df_name], row[minor_col_in_df_name], \
                        row[product_name_col_in_df])
                    continue
                
                pi = rp.concept_rt_to_ProgramInformation(major_value=major_key, minor_value=minor_key,\
                    name=row[product_name_col_in_df], my_row=row, roadmap_configuration=roadmap_configuration,\
                     from_roadmap_top_cm=y_offset + (row_count * PROG_PROG_VERT_DISTANCE_CM),\
                        missing_ms=NOT_PROVIDED)
                print('-----', pi.annotations())
                
                ### Note - the render program function actually updates each of the
                ### milestones with x axis information about where they were rendered.
                pi = roadmap_slide.render_program(program_information=pi, align_zero=align_zero)               
                
                shorthand_name_dict = add_to_ordered_dict(shorthand_name_dict,major_key,minor_key,pi.shorthand_name, pi)

                row_count = row_count + 1
        
    return shorthand_name_dict

"""
def render_annotations( shorthand_name_dict, roadmap_slide, roadmap_canvas, roadmap_table,\
    align_zero=False ):
    if roadmap_table.annotations_list() is not None:
        for annoation in roadmap_table.annotations_list():
            if annoation[rt.TYPE] == rt.ANNOTTION_SLIP:
                try:
                    start_pi = shorthand_name_dict[annoation[rt.SLIP_START_SHORT_NAME]]
                except:
                    print("Annotate: Cannot find:", annoation[rt.SLIP_START_SHORT_NAME] )
                    continue
                try:
                    start_ms = start_pi.milestone_by_name(milestone_name=annoation[rt.SLIP_START_MILESTONE])
                except:
                    print("Annotate: Cannot find:", annoation[rt.SLIP_START_MILESTONE] )
                    continue
                try:
                    end_pi = shorthand_name_dict[annoation[rt.SLIP_END_SHORT_NAME]]
                except:
                    print("Annotate: Cannot find:", annoation[rt.SLIP_END_SHORT_NAME] )
                    continue
                try:
                    end_ms = end_pi.milestone_by_name(milestone_name=annoation[rt.SLIP_END_MILESTONE])
                except:
                    print("Annotate: Cannot find:", annoation[rt.SLIP_END_MILESTONE] )
                    continue
                
                print(start_pi, start_ms, end_pi, end_ms)

                roadmap_slide.annotate_slip(start_pi=start_pi, start_ms=start_ms, end_pi=end_pi,\
                    end_ms=end_ms, align_zero=align_zero)
            
            if annoation[rt.TYPE] == rt.ANNOTATION_INLINEWWS:
                for shorthand_name, program_information in shorthand_name_dict.items():
                    roadmap_slide.annotate_inline(pi=program_information)
        
    for shorthand_name, program_information in shorthand_name_dict.items():
            roadmap_slide.annotate_backend(pi=program_information,a0_ti_text='A0TI', prq_text='PRQ')

    return
            
"""
def render_annotations( shorthand_name_dict, roadmap_slide, roadmap_canvas, align_zero=False ):
    for key1, major in shorthand_name_dict.items():
        for key2, minor in major.items():
            for key3, pi in minor.items():
                print('>>>>>>>>', key3, pi.annotations())
                roadmap_slide.render_annotations(pi=pi, all_pis_dict=shorthand_name_dict)
    
def render_roadmap_from_paths(golden_doc_path, roadmap_helper_path, roadmap_config_path, \
    roadmap_template_path, start_ww, end_ww, roadmap_title=None, \
        roadmap_top_cm=1.5, roadmap_height_cm=16.5, input_slide_index=0, \
            align_zero=False ):
    
    if golden_doc_path is not None:
        ### Read roadmap information from "golden doc" and helper files
        roadmap_table = rt.RoadmapTable(path_to_roadmap=golden_doc_path, \
            path_to_helper_file=roadmap_helper_path)
    else:
        raise ValueError("Must provide valid golden_doc_path")

    
    print(roadmap_table.annotations_list())
    
    #roadmap_table.df.to_excel('/Users/scotttan/Intel Corporation/Graphics Golden Doc Development - RoadmapPPTGeneration/ElastiScenarios/OutputAfterFixUp.xlsx')
    
    ### Read in the roadmap input parsing and output configuration information
    r_c = rc.RoadmapConfig(roadmap_config_path)
    print(r_c.swimlanes_hierarchy)

    ### Create the rendering canvas 
    cag = rp.CanvasGrid(start_ww=start_ww,end_ww=end_ww)
    cac = rp.RoadmapCanvas(grid=cag, canvas_top_cm=roadmap_top_cm, canvas_height_cm=roadmap_height_cm)

    ### Create the swimlane grid and table
    slg = rp.SwimlanesGrid(hierarchy=r_c.swimlanes_hierarchy, \
        major_column_name=r_c.major_column_name,\
        minor_column_name=r_c.minor_column_name)
        
    slt = rp.SwimlaneTable(grid=slg,table_top_cm=cac.top_of_draw_area,\
        table_left_edge_cm=0.0, table_height_cm=cac.draw_area_height)

    ## Read in the presentation
    pptx = Presentation(roadmap_template_path)

    ## Create the RoadmapSlide object so we can render onto it
    rs = rp.RoadmapSlide(slide=pptx.slides[input_slide_index],roadmap_canvas=cac,\
        roadmap_top_cm=roadmap_top_cm)

    render_swimlanes_and_qts(roadmap_slide=rs, roadmap_canvas=cac, swimlane_table=slt)

    shorthand_name_dict= render_roadmap(roadmap_slide=rs, roadmap_configuration=r_c,\
         roadmap_table=roadmap_table, swimlane_table=slt,align_zero=align_zero)
        
    #for key, value in shorthand_name_dict.items():
    #    print('-'*2,key)
    #    for ms in value.milestones:
    #        print('-'*4,f"{ms.marker_left_cm:f2.2}", end='')
    #        print('-'*4,f"{ms.marker_width_cm:f2.2}")
    
    render_annotations(shorthand_name_dict=shorthand_name_dict,roadmap_slide=rs,\
        roadmap_canvas=r_c, align_zero=align_zero)

    return pptx

if __name__ == '__main__':
    import pandas as pd
    from pptx import Presentation
    from pptx.util import Inches, Pt, Cm
    from collections import  OrderedDict

    #### FILES #####
    GOLDEN_DOC_XLS_PATH = \
    "/Users/scotttan/Intel Corporation/Graphics Golden Doc Development - RoadmapPPTGeneration/ElastiScenarios/ElastiScenariosGoldenDoc.xlsx"
    ROADMAP_HELPER_EXCEL = \
            '/Users/scotttan/Intel Corporation/Graphics Golden Doc Development - RoadmapPPTGeneration/ElastiScenarios/ElastiScenariosHelperDoc.xlsx'
    ROADMAP_CONFIG_EXCEL = '/Users/scotttan/Intel Corporation/Graphics Golden Doc Development - RoadmapPPTGeneration/ElastiScenarios/ElastiConfig.xlsx'
    ROADMAP_TEMPLATE = '/Users/scotttan/Intel Corporation/Graphics Golden Doc Development - RoadmapPPTGeneration/ScottsWorkbench/BlankIAGSTemplate.pptx'
    ROADMAP_OUT = '/Users/scotttan/Intel Corporation/Graphics Golden Doc Development - RoadmapPPTGeneration/ElastiScenarios/IAGS_OUT.pptx'

    ### Constant configuraiton - will need to be paramters
    START_WW = iw.WW(25,20)
    END_WW = iw.WW(51,23)
    ROADMAP_TOP_CM = 1.5
    ROADMAP_HEIGHT_CM =16.25
    INPUT_SLIDE_INDEX = 0

    SIMPLE_ROADMAP = None
    CONCEPT_ROADMAP = '/Users/scotttan/Intel Corporation/Graphics Golden Doc Development - RoadmapPPTGeneration/ElastiScenarios/PreConceptBronzeDoc.xlsx'
    
    pptx = render_roadmap_from_paths(
        golden_doc_path=SIMPLE_ROADMAP,
        concept_doc_path=CONCEPT_ROADMAP,
        roadmap_helper_path=None,\
        roadmap_config_path=ROADMAP_CONFIG_EXCEL, 
        roadmap_template_path=ROADMAP_TEMPLATE, \
        start_ww=START_WW, end_ww=END_WW,\
        roadmap_top_cm=ROADMAP_TOP_CM, roadmap_height_cm=ROADMAP_HEIGHT_CM,\
        input_slide_index=INPUT_SLIDE_INDEX)
    
    pptx.save(ROADMAP_OUT)







    